#!/usr/bin/env python3
"""Rediseno de las piezas comerciales de Tazki (pedido de Diego, 17/18-ago-2026).

Uso:  python3 ventas_rediseno.py [--solo construccion]

Rehace tres laminas de cada presentacion de ventas-piezas/, con el sistema de diseno
del deck (DM Sans, navy 0C1E50, cobalto 1537D1, lienzo 13.33 x 7.5 pulgadas):

  1. PORTADA   -> foto de prevencion nueva por industria (ventas-assets/seleccion-fotos.json)
  2. DS 44     -> lamina de contexto con foto a sangre, titular grande (57pt) y 3 puntos
  3. LOGOS     -> muro de logos de clientes de esa industria (ventas-assets/logos-clientes/,
                  clasificados en ventas-assets/logos-por-industria.json)

NO toca las laminas de modulos: ahi viven los GIF de los flujos, que se conservan.
Es idempotente: la lamina DS 44 se reconstruye desde cero en cada corrida.
"""
import glob, json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

RAIZ    = os.path.dirname(os.path.abspath(__file__))
PREV    = os.path.join(RAIZ, "ventas-assets/prev-nuevas")
LOGOS   = os.path.join(RAIZ, "ventas-assets/logos-recortados")
VACIAS  = set(json.load(open(os.path.join(os.path.dirname(os.path.abspath(__file__)),
          "ventas-assets/logos-recorte-reporte.json"), encoding="utf-8"))["tarjeta_vacia"])
SEL     = json.load(open(os.path.join(RAIZ, "ventas-assets/seleccion-fotos.json"), encoding="utf-8"))
MAPA    = json.load(open(os.path.join(RAIZ, "ventas-assets/logos-por-industria.json"), encoding="utf-8"))

NAVY, COBALT = RGBColor(0x0C,0x1E,0x50), RGBColor(0x15,0x37,0xD1)
GRIS, BLANCO = RGBColor(0x5B,0x64,0x72), RGBColor(0xFF,0xFF,0xFF)
CELESTE, AZUL_CLARO = RGBColor(0xC9,0xD3,0xEA), RGBColor(0x7F,0xB2,0xFF)
LINEA = RGBColor(0xE2,0xE7,0xF1)
FUENTE, MARCA = "DM Sans", "​"
W, H = 13.333, 7.5

# carpeta de la pieza -> etiqueta que se muestra en la lamina de logos
ETIQUETA = {
 "construccion":"Construcción · obras · contratistas", "mineria":"Minería y faenas",
 "transporte":"Transporte, logística y almacenamiento", "manufactura":"Industria y manufactura",
 "energia":"Energía, sanitarias y servicios básicos", "comercio":"Comercio y retail",
 "salud":"Salud y asistencia social", "agro":"Agro y acuicultura",
 "telecom":"Telecom, TI y servicios", "contratistas":"Gestión de contratistas",
 "core":"Empresas de todos los rubros", "expandido":"Empresas de todos los rubros",
 "enterprise":"Operaciones multisitio", "plan-core":"Empresas de todos los rubros",
}

# ---------------------------------------------------------------- utilidades

def _txt(s, x, y, w, h, texto, size, bold, color, line=1.15, align=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, l in enumerate(texto.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line
        if align: p.alignment = align
        r = p.add_run(); r.text = l
        r.font.name = FUENTE; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def _rect(s, x, y, w, h, fill, forma=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(forma, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if radius is not None:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return sh

def _foto_cubriendo(s, ruta, x, y, w, h):
    """Inserta la foto recortada tipo 'cover' en el rectangulo dado."""
    iw, ih = Image.open(ruta).size
    escala = max(w / (iw/96), h / (ih/96))
    aw, ah = (iw/96)*escala, (ih/96)*escala
    pic = s.shapes.add_picture(ruta, Inches(x - (aw-w)/2), Inches(y - (ah-h)/2), Inches(aw), Inches(ah))
    # recorte al marco
    pic.crop_left = pic.crop_right = max(0.0, (aw-w)/aw/2)
    pic.crop_top = pic.crop_bottom = max(0.0, (ah-h)/ah/2)
    pic.left, pic.top = Inches(x), Inches(y)
    pic.width, pic.height = Inches(w), Inches(h)
    return pic

def _velo(s, x, y, w, h, color, alpha_pct):
    """Rectangulo semitransparente sobre la foto (PowerPoint: transparencia del relleno)."""
    sh = _rect(s, x, y, w, h, color)
    from lxml import etree
    sp = sh.fill._xPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    clr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    a = etree.SubElement(clr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
    a.set('val', str(int((100-alpha_pct)*1000)))
    return sh

def _vaciar(s):
    """Deja la lamina en blanco conservando su lugar en el mazo."""
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def _nueva(prs):
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders._element))
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s

def _mover(prs, desde_final_a):
    lst = prs.slides._sldIdLst; ids = list(lst)
    lst.remove(ids[-1]); lst.insert(desde_final_a, ids[-1])

def _borrar(prs, idx):
    lst = prs.slides._sldIdLst; ids = list(lst)
    rid = ids[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rid); lst.remove(ids[idx])

# ---------------------------------------------------------------- laminas

def lamina_ds44(s, foto):
    _txt(s, 0.05, 7.42, 0.3, 0.06, MARCA, 1, False, BLANCO)
    _foto_cubriendo(s, foto, 0, 0, W, H)
    _velo(s, 0, 0, W, H, NAVY, 82)
    _velo(s, 0, 0, 7.4, H, NAVY, 94)

    _rect(s, 0.77, 0.70, 0.36, 0.035, AZUL_CLARO)
    _txt(s, 1.25, 0.60, 6.0, 0.3, "EL CONTEXTO", 13, True, AZUL_CLARO)
    _txt(s, 0.77, 1.02, 7.0, 2.4, "Ya no basta\ncon tener\nlos papeles.", 54, True, BLANCO, line=1.02)
    _txt(s, 0.80, 3.42, 5.9, 1.0,
         "El DS 44 cambió la regla: desde 2025 la prevención dejó de medirse por documentos "
         "archivados y pasó a medirse por gestión demostrable.", 15.5, False, CELESTE, line=1.42)

    puntos = [("01","La ODI se convirtió en IRL",
               "Informar los riesgos de cada puesto, con firma de quien los conoció, y actualizarlos cuando la tarea cambia."),
              ("02","La MIPER manda el programa",
               "La matriz define el programa preventivo, y ese programa vence a los 30 días de actualizar la matriz."),
              ("03","Se fiscaliza la evidencia",
               "El FUF no revisa si el documento existe: revisa fecha, responsable, firma y seguimiento.")]
    for i, (n, tit, cue) in enumerate(puntos):
        x = 0.77 + i*4.0
        _rect(s, x, 4.72, 3.55, 0.04, COBALT)
        _txt(s, x, 4.88, 3.5, 0.25, n, 13, True, AZUL_CLARO)
        _txt(s, x, 5.18, 3.5, 0.6, tit, 18, True, BLANCO, line=1.12)
        _txt(s, x, 5.72, 3.5, 1.0, cue, 11.5, False, CELESTE, line=1.36)

    _rect(s, 0, 6.85, W, 0.65, COBALT)
    _txt(s, 0.77, 7.03, 11.8, 0.35,
         "Tazki es donde esa gestión queda registrada, al día y lista para mostrar.", 15, True, BLANCO)
    return s

def lamina_logos(s, carpeta, paginas):
    _txt(s, 0.05, 7.42, 0.3, 0.06, MARCA, 1, False, BLANCO)
    _rect(s, 0, 0, W, H, BLANCO)
    _rect(s, 0, 0, 0.16, H, COBALT)

    _txt(s, 0.77, 0.62, 11.5, 0.5, "Empresas que ya confían en Tazki", 31, True, NAVY)
    _txt(s, 0.79, 1.24, 11.0, 0.35,
         f"{ETIQUETA.get(carpeta,'Empresas de todos los rubros')} · {len(paginas)} de las +150 que ya digitalizaron su prevención.",
         14, False, GRIS)

    archivos = [os.path.join(LOGOS, f"{p}.png") for p in paginas if p not in VACIAS]
    archivos = [a for a in archivos if os.path.exists(a)]
    n = len(archivos)
    # se elige el numero de columnas que deja la ultima fila mas llena (evita filas de 1 logo)
    def _huecos(c):
        f = (n + c - 1)//c
        return (f*c - n, abs(c - 7))          # a igualdad de huecos, se prefiere ~7 columnas
    cols = min(range(4, 9), key=_huecos)
    filas = (n + cols - 1)//cols
    cw, gap = (W - 1.54 - (cols-1)*0.14)/cols, 0.14
    ch = min(1.16, (5.55 - (filas-1)*0.14)/max(filas,1))
    y0 = 1.95 + max(0, (5.35 - (filas*ch + (filas-1)*0.14))/2)

    for i, a in enumerate(archivos):
        cx = 0.77 + (i % cols)*(cw+gap)
        cy = y0 + (i//cols)*(ch+gap)
        iw, ih = Image.open(a).size
        cajaw, cajah = cw*0.99, ch*0.99
        esc = min(cajaw/(iw/96), cajah/(ih/96))
        aw, ah = (iw/96)*esc, (ih/96)*esc
        s.shapes.add_picture(a, Inches(cx+(cw-aw)/2), Inches(cy+(ch-ah)/2), Inches(aw), Inches(ah))
    return s

# ---------------------------------------------------------------- proceso

def procesar(carpeta):
    ruta = glob.glob(f"{RAIZ}/ventas-piezas/{carpeta}/*.pptx")
    if not ruta: return False
    ruta = ruta[0]
    prs = Presentation(ruta)

    # paginas de logos de esta industria
    paginas = sorted(int(k) for k,v in MAPA["logos"].items() if v["pieza"] == carpeta)
    if not paginas:   # piezas sin industria propia: se usa una muestra transversal
        paginas = sorted(int(k) for k in MAPA["logos"])[:28]

    foto_portada = os.path.join(PREV, SEL["portadas"].get(carpeta, SEL["portadas"]["core"]))
    foto_ds44    = os.path.join(PREV, SEL["ds44_todas"])

    # 1) portada: se reemplaza la foto de fondo (la imagen que cubre toda la lamina)
    p0 = list(prs.slides)[0]
    for sh in list(p0.shapes):
        if sh.shape_type == 13 and sh.width > Inches(12):
            sh._element.getparent().remove(sh._element)
            pic = _foto_cubriendo(p0, foto_portada, 0, 0, W, H)
            p0.shapes._spTree.remove(pic._element)
            p0.shapes._spTree.insert(2, pic._element)   # al fondo
            _velo(p0, 0, 0, W, H, NAVY, 68)
            velo = p0.shapes[-1]
            p0.shapes._spTree.remove(velo._element)
            p0.shapes._spTree.insert(3, velo._element)
            break

    # 2) y 3): se vacian y se redibujan en su mismo lugar (borrar laminas deja partes huerfanas)
    laminas = list(prs.slides)
    lamina_ds44(_vaciar(laminas[1]), foto_ds44)
    lamina_logos(_vaciar(laminas[2]), carpeta, paginas)

    prs.save(ruta)
    print(f"  ✓ {carpeta:<14} portada + DS44 + {len(paginas):>2} logos   ({len(prs.slides)} láminas)")
    return True

if __name__ == "__main__":
    solo = sys.argv[2] if len(sys.argv) > 2 and sys.argv[1] == "--solo" else None
    carpetas = [solo] if solo else sorted(os.path.basename(d) for d in glob.glob(f"{RAIZ}/ventas-piezas/*") if os.path.isdir(d))
    print(f"Rediseñando {len(carpetas)} pieza(s)…")
    for c in carpetas: procesar(c)
