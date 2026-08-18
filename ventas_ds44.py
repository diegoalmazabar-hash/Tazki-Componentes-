#!/usr/bin/env python3
"""Inserta la lamina de contexto DS 44 al inicio de las presentaciones comerciales.

Uso:  python3 ventas_ds44.py <pptx> [pptx...]
      python3 ventas_ds44.py --todas

Pedido de Diego (17-ago-2026): las presentaciones deben partir explicando que el DS 44
cambio la gestion preventiva y que Tazki ayuda con eso, antes de mostrar producto.

La lamina se inserta como pagina 2 (despues de la portada, antes del muro de logos) y
respeta el sistema de diseno del deck: DM Sans, navy 0C1E50, cobalto 1537D1,
tarjetas F4F6FB, lienzo 13.33 x 7.5 pulgadas.

Es idempotente: si la lamina ya existe (se marca con MARCA), no la vuelve a insertar.
"""
import copy, glob, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY   = RGBColor(0x0C, 0x1E, 0x50)
COBALT = RGBColor(0x15, 0x37, 0xD1)
GRIS   = RGBColor(0x5B, 0x64, 0x72)
CARD   = RGBColor(0xF4, 0xF6, 0xFB)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CELESTE= RGBColor(0xCF, 0xE0, 0xFA)
FUENTE = "DM Sans"
MARCA  = "​"  # zero-width space: marca invisible para no duplicar la lamina

KICKER = "EL CONTEXTO"
TITULO = "El DS 44 cambió lo que se te exige:\nya no basta con tener los papeles."
BAJADA = ("Desde 2025 la prevención dejó de medirse por documentos archivados y pasó a medirse "
          "por gestión demostrable. Estas son las tres exigencias que cambiaron el día a día.")

TARJETAS = [
    ("La ODI se convirtió en IRL",
     "Ahora hay que informar los riesgos de cada puesto de trabajo, dejar firma de que la persona "
     "los conoció y actualizar esa información cada vez que el riesgo o la tarea cambian."),
    ("La MIPER manda el programa",
     "La matriz de peligros define el programa de trabajo preventivo, y ese programa debe elaborarse "
     "o modificarse dentro de los 30 días siguientes a confeccionar o actualizar la matriz."),
    ("La fiscalización pide evidencia",
     "El Formulario Único de Fiscalización no revisa si el documento existe: revisa fecha, responsable, "
     "firma y seguimiento. Sin trazabilidad, no hay cómo demostrar la gestión."),
]
CIERRE = "Tazki es donde esa gestión queda registrada, al día y lista para mostrar."


def _txt(slide, x, y, w, h, texto, size, bold, color, space_after=0, line=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linea in enumerate(texto.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = linea
        r.font.name = FUENTE; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def _card(slide, x, y, w, h, fill, radius=0.045):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.text_frame.text = ""
    return sh


def construir(prs):
    # lienzo en blanco: se usa el layout mas vacio disponible
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders._element))
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)

    # marca invisible para idempotencia
    _txt(s, 0.05, 7.42, 0.3, 0.06, MARCA, 1, False, BLANCO)

    _txt(s, 0.9, 0.72, 9.0, 0.3, KICKER, 12.5, True, COBALT)
    _txt(s, 0.9, 1.05, 11.2, 1.2, TITULO, 30, True, NAVY, line=1.18)
    _txt(s, 0.92, 2.18, 10.6, 0.7, BAJADA, 14.5, False, GRIS, line=1.35)

    xs = [0.9, 4.85, 8.8]
    for x, (tit, cuerpo) in zip(xs, TARJETAS):
        _card(s, x, 3.15, 3.75, 2.75, CARD)
        num_bg = _card(s, x + 0.35, 3.5, 0.5, 0.5, BLANCO, radius=0.5)
        _txt(s, x + 0.35, 3.6, 0.5, 0.3, str(xs.index(x) + 1), 15, True, COBALT).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _txt(s, x + 0.35, 4.2, 3.05, 0.7, tit, 15.5, True, NAVY, line=1.2)
        _txt(s, x + 0.37, 4.82, 3.03, 1.5, cuerpo, 11.0, False, GRIS, line=1.35)

    barra = slide_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.33), Inches(1.0))
    barra.fill.solid(); barra.fill.fore_color.rgb = COBALT
    barra.line.fill.background(); barra.shadow.inherit = False
    _txt(s, 0.9, 6.83, 11.5, 0.4, CIERRE, 16, True, BLANCO)
    return s


def ya_tiene(prs):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and MARCA in sh.text_frame.text:
                return True
    return False


def mover_a_posicion(prs, indice_destino):
    """python-pptx solo agrega al final: se reordena el sldIdLst."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[-1])
    lst.insert(indice_destino, ids[-1])


def procesar(ruta):
    prs = Presentation(ruta)
    if ya_tiene(prs):
        print(f"  = ya tenía la lámina, sin cambios: {ruta}")
        return False
    antes = len(prs.slides)
    construir(prs)
    mover_a_posicion(prs, 1)          # queda como lámina 2
    prs.save(ruta)
    print(f"  + DS 44 insertada como lámina 2 ({antes} → {antes+1}): {ruta}")
    return True


if __name__ == "__main__":
    args = sys.argv[1:]
    rutas = sorted(glob.glob("ventas-piezas/*/*.pptx")) if (not args or args[0] == "--todas") else args
    print(f"Procesando {len(rutas)} presentación(es)…")
    n = sum(procesar(r) for r in rutas)
    print(f"Listo: {n} modificada(s).")
