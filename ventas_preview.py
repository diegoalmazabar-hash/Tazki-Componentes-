#!/usr/bin/env python3
"""Simula con PIL como se ve una lamina del PPTX, para poder revisarla sin PowerPoint.

Uso:  python3 ventas_preview.py <carpeta> <n_lamina> <salida.png>

Lee la lamina real del pptx (posiciones, textos, tamanos, colores, imagenes, velos con
transparencia) y la dibuja. No es un render perfecto —la tipografia es la del sistema y
no DM Sans— pero sirve para pillar solapes, textos ilegibles y velos mal aplicados,
que es justo lo que no se puede ver leyendo el XML.
"""
import sys, os, io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

K = 96  # pixeles por pulgada
NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def fuente(px, bold):
    rutas = ["/usr/share/fonts/truetype/dejavu/DejaVuSans%s.ttf" % ("-Bold" if bold else ""),
             "/usr/share/fonts/truetype/liberation/LiberationSans%s.ttf" % ("-Bold" if bold else "")]
    for r in rutas:
        if os.path.exists(r):
            try: return ImageFont.truetype(r, max(8, int(px)))
            except Exception: pass
    return ImageFont.load_default()

def emu_px(v):
    return int(Emu(v).inches * K)

def opacidad(shape):
    """Devuelve la opacidad del relleno (0-1); 1 si es solido."""
    try:
        sf = shape.fill._xPr.find(NS + 'solidFill')
        if sf is None: return 1.0
        clr = sf.find(NS + 'srgbClr')
        if clr is None: return 1.0
        a = clr.find(NS + 'alpha')
        return int(a.get('val'))/100000.0 if a is not None else 1.0
    except Exception:
        return 1.0

def render(pptx, idx, salida):
    prs = Presentation(pptx)
    s = list(prs.slides)[idx-1]
    W, H = int(Emu(prs.slide_width).inches*K), int(Emu(prs.slide_height).inches*K)
    lienzo = Image.new("RGB", (W, H), "white")

    for sh in s.shapes:
        x, y = emu_px(sh.left or 0), emu_px(sh.top or 0)
        w, h = max(1, emu_px(sh.width or 0)), max(1, emu_px(sh.height or 0))

        if sh.shape_type == 13:  # imagen
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
            except Exception:
                continue
            cl, cr = getattr(sh, "crop_left", 0) or 0, getattr(sh, "crop_right", 0) or 0
            ct, cb = getattr(sh, "crop_top", 0) or 0, getattr(sh, "crop_bottom", 0) or 0
            if cl or cr or ct or cb:
                iw, ih = im.size
                im = im.crop((int(iw*cl), int(ih*ct), int(iw*(1-cr)), int(ih*(1-cb))))
            lienzo.paste(im.resize((w, h)), (x, y))
            continue

        # relleno (con transparencia si la tiene)
        try:
            if sh.fill.type == 1:
                col = sh.fill.fore_color.rgb
                capa = Image.new("RGB", (w, h), (col[0], col[1], col[2]) if not isinstance(col, str) else "#"+str(col))
                op = opacidad(sh)
                region = lienzo.crop((x, y, x+w, y+h))
                if region.size == capa.size:
                    lienzo.paste(Image.blend(region, capa, min(1.0, max(0.0, op))), (x, y))
        except Exception:
            pass

        # texto
        if sh.has_text_frame and sh.text_frame.text.strip():
            d = ImageDraw.Draw(lienzo)
            cy = y
            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs)
                if not txt.strip():
                    cy += 10; continue
                r0 = p.runs[0]
                pt = r0.font.size.pt if r0.font.size else 14
                bold = bool(r0.font.bold)
                try: col = r0.font.color.rgb; col = (col[0], col[1], col[2])
                except Exception: col = (0, 0, 0)
                f = fuente(pt*K/72.0, bold)
                ls = p.line_spacing or 1.15
                # ajuste de linea simple al ancho de la caja
                palabras, linea = txt.split(" "), ""
                for pal in palabras:
                    prueba = (linea+" "+pal).strip()
                    if d.textlength(prueba, font=f) <= w or not linea:
                        linea = prueba
                    else:
                        d.text((x, cy), linea, font=f, fill=col); cy += int(pt*K/72.0*ls); linea = pal
                if linea:
                    d.text((x, cy), linea, font=f, fill=col); cy += int(pt*K/72.0*ls)

    lienzo.save(salida)
    print(f"{salida}  ({W}x{H})  lámina {idx} de {os.path.basename(pptx)}")

if __name__ == "__main__":
    render(sys.argv[1], int(sys.argv[2]), sys.argv[3])
